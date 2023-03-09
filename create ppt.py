import collections 
import collections.abc
import pptx
from pptx.util import Inches
from pptx.chart.data import CategoryChartData
from datetime import datetime as dt
import pandas as pd
import functions as fci

df = pd.read_excel("Gastos.xlsx")
img = 'imagens/background.png'

ppt = pptx.Presentation()
slide1 = ppt.slides.add_slide(ppt.slide_layouts[6])

left = top = Inches(0)


#############################################################################
#Primeiro slide 

#add fundo
pic = slide1.shapes.add_picture(img, left, top, width=Inches(10), height=Inches(7.5))

#add logo
#pic = slide1.shapes.add_picture('imagens/Python.png', Inches(0.2), Inches(0.2))


#Text box (x, y, altura, largura)
textbox = slide1.shapes.add_textbox(Inches(6), Inches(3.5), Inches(4), Inches(3))

    #titulo
titulo = textbox.text_frame.add_paragraph()
titulo.text = "Angelo Queiroz"
titulo.font.color.rgb = pptx.dml.color.RGBColor(255, 255, 255)
titulo.font.size = pptx.util.Pt(44)
#titulo.font.bold = True

    #subtitulo
text_frame = textbox.text_frame
subtitulo = text_frame.add_paragraph()
subtitulo = text_frame.add_paragraph()
subtitulo = text_frame.add_paragraph()
subtitulo.text = f"{dt.today().strftime('%d.%m.%Y')}"
subtitulo.font.color.rgb = pptx.dml.color.RGBColor(255, 255, 255)
subtitulo.font.size = pptx.util.Pt(28)
textbox.text_frame.paragraphs[4].alignment = pptx.enum.text.PP_ALIGN.RIGHT

#############################################################################
#############################################################################
#Segundo slide

slide1 = ppt.slides.add_slide(ppt.slide_layouts[6])

    #Titulo
textbox = slide1.shapes.add_textbox(Inches(0.1), Inches(-0.3), Inches(1), Inches(1))
paragrafo = textbox.text_frame.add_paragraph()
paragrafo.text = 'Gastos'
paragrafo.font.color.rgb = pptx.dml.color.RGBColor(0, 102, 255)
paragrafo.font.size = pptx.util.Pt(32)
paragrafo.font.bold = True

    #Gastos totais
textbox = slide1.shapes.add_textbox(Inches(1.3), Inches(5.5), Inches(1.7), Inches(1))
paragrafo = textbox.text_frame.add_paragraph()
paragrafo.text = 'Total'
paragrafo.font.color.rgb = pptx.dml.color.RGBColor(0, 102, 255)
paragrafo.font.size = pptx.util.Pt(22)
paragrafo.font.bold = False

paragrafo = textbox.text_frame.add_paragraph()
gastos = df['Gasto'].sum()
paragrafo.text = f'R$:{gastos}'
paragrafo.font.color.rgb = pptx.dml.color.RGBColor(0, 0, 0)
paragrafo.font.size = pptx.util.Pt(18)
paragrafo.font.bold = False

    # Periodo
textbox = slide1.shapes.add_textbox(Inches(3), Inches(5.5), Inches(3.3), Inches(1))
paragrafo = textbox.text_frame.add_paragraph()
paragrafo.text = 'Período'
paragrafo.font.color.rgb = pptx.dml.color.RGBColor(0, 102, 255)
paragrafo.font.size = pptx.util.Pt(22)
paragrafo.font.bold = False

paragrafo = textbox.text_frame.add_paragraph()
dia_min = df['Dia'].min().strftime('%d/%m/%Y')
dia_max = df['Dia'].max().strftime('%d/%m/%Y')
paragrafo.text = f'{dia_min} - {dia_max}'
paragrafo.font.color.rgb = pptx.dml.color.RGBColor(0, 0, 0)
paragrafo.font.size = pptx.util.Pt(18)
paragrafo.font.bold = False
textbox.text_frame.paragraphs[1].alignment = pptx.enum.text.PP_ALIGN.CENTER
textbox.text_frame.paragraphs[2].alignment = pptx.enum.text.PP_ALIGN.CENTER

    # MAIOR GASTO

textbox = slide1.shapes.add_textbox(Inches(5), Inches(5.5), Inches(5.4), Inches(5.6))
paragrafo = textbox.text_frame.add_paragraph()
paragrafo.text = 'Maior gasto'
paragrafo.font.color.rgb = pptx.dml.color.RGBColor(0, 102, 255)
paragrafo.font.size = pptx.util.Pt(22)
paragrafo.font.bold = False

paragrafo = textbox.text_frame.add_paragraph()
valor, categoria = fci.maior_gastos(df)
paragrafo.text = f'{categoria} - {valor}'
paragrafo.font.color.rgb = pptx.dml.color.RGBColor(0, 0, 0)
paragrafo.font.size = pptx.util.Pt(18)
paragrafo.font.bold = False
textbox.text_frame.paragraphs[1].alignment = pptx.enum.text.PP_ALIGN.CENTER
textbox.text_frame.paragraphs[2].alignment = pptx.enum.text.PP_ALIGN.CENTER

    # Graficos

valores, categoria = fci.gastos_por_categoria(df)
  #y        x

dados_grafico = CategoryChartData()
dados_grafico.categories = categoria
dados_grafico.add_series("Valores", valores)
# tipos de gráfico: https://python-pptx.readthedocs.io/en/latest/api/enum/XlChartType.html#xlcharttype
slide1.shapes.add_chart(pptx.enum.chart.XL_CHART_TYPE.COLUMN_CLUSTERED, Inches(1), Inches(1), Inches(7.5), Inches(4.5), dados_grafico)

#############################################################################

ppt.save("PPT.pptx")
print("Apresentação criada com sucesso!")
